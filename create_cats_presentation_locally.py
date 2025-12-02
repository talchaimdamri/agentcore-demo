from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a presentation object
prs = Presentation()

# Set slide dimensions (standard 16:9)
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "All About Cats"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(54)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.8))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Our Feline Friends"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(28)
subtitle_para.alignment = PP_ALIGN.CENTER

# Slide 2: Introduction to Cats
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title2_frame = title2.text_frame
title2_frame.text = "What Are Cats?"
title2_para = title2_frame.paragraphs[0]
title2_para.font.size = Pt(40)
title2_para.font.bold = True

content2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5))
content2_frame = content2.text_frame
content2_frame.word_wrap = True

points = [
    "Domesticated carnivorous mammals",
    "Scientific name: Felis catus",
    "Estimated 600 million cats worldwide",
    "Domesticated around 10,000 years ago",
    "Independent yet affectionate companions"
]

for point in points:
    p = content2_frame.add_paragraph()
    p.text = "• " + point
    p.font.size = Pt(24)
    p.space_before = Pt(12)

# Slide 3: Cat Behaviors
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title3_frame = title3.text_frame
title3_frame.text = "Fascinating Cat Behaviors"
title3_para = title3_frame.paragraphs[0]
title3_para.font.size = Pt(40)
title3_para.font.bold = True

content3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5))
content3_frame = content3.text_frame
content3_frame.word_wrap = True

behaviors = [
    "Purring: Sign of contentment (or sometimes stress)",
    "Kneading: Behavior from kittenhood",
    "Grooming: Spend 30-50% of day cleaning themselves",
    "Sleeping: Sleep 12-16 hours per day",
    "Hunting instinct: Natural predators, even when domesticated"
]

for behavior in behaviors:
    p = content3_frame.add_paragraph()
    p.text = "• " + behavior
    p.font.size = Pt(22)
    p.space_before = Pt(12)

# Slide 4: Fun Facts
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title4_frame = title4.text_frame
title4_frame.text = "Fun Cat Facts! 🐱"
title4_para = title4_frame.paragraphs[0]
title4_para.font.size = Pt(40)
title4_para.font.bold = True

content4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5))
content4_frame = content4.text_frame
content4_frame.word_wrap = True

facts = [
    "Cats have 32 muscles in each ear",
    "A group of cats is called a 'clowder'",
    "Cats can rotate their ears 180 degrees",
    "They have a third eyelid for protection",
    "Cats spend 70% of their lives sleeping"
]

for fact in facts:
    p = content4_frame.add_paragraph()
    p.text = "✓ " + fact
    p.font.size = Pt(24)
    p.space_before = Pt(14)

# Slide 5: Closing
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
closing_box = slide5.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
closing_frame = closing_box.text_frame
closing_frame.text = "Thank You!"
closing_para = closing_frame.paragraphs[0]
closing_para.font.size = Pt(54)
closing_para.font.bold = True
closing_para.alignment = PP_ALIGN.CENTER

tagline_box = slide5.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.8))
tagline_frame = tagline_box.text_frame
tagline_frame.text = "Cats: Purr-fect Companions! 🐾"
tagline_para = tagline_frame.paragraphs[0]
tagline_para.font.size = Pt(28)
tagline_para.alignment = PP_ALIGN.CENTER

# Save the presentation
prs.save('cats_presentation.pptx')
print("✅ Presentation created successfully!")
print("📄 File saved as: cats_presentation.pptx")
print(f"📊 Total slides: {len(prs.slides)}")
